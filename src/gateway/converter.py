from __future__ import annotations

import hashlib
import json
import logging
from dataclasses import asdict, dataclass
from pathlib import Path

from document_image_renderer import (
    RenderedImage,
    RendererError,
    RenderOptions,
    UnsupportedFormatError,
    render_document,
)


SUPPORTED_EXTENSIONS = {".pdf", ".pptx", ".docx", ".xlsx"}
logger = logging.getLogger("uvicorn.error")


@dataclass(frozen=True)
class Artifact:
    page_number: int
    image_path: str
    text_path: str
    width: int
    height: int
    media_type: str
    sha256: str


@dataclass(frozen=True)
class ConversionResult:
    artifacts: list[Artifact]
    warnings: list[str]


class DocumentPageLimitError(RendererError):
    pass


def convert_document(source: Path, output_dir: Path, max_pages: int = 20) -> ConversionResult:
    extension = source.suffix.lower()
    if extension not in SUPPORTED_EXTENSIONS:
        raise UnsupportedFormatError(f"Unsupported file type: {extension or '(none)'}")
    logger.info("Document rendering started: filename=%s format=%s", source.name, extension)
    text_blocks = extract_text(source, extension)
    render_result = render_document(
        source,
        output_dir,
        options=RenderOptions(
            image_format="png",
            dpi=150,
            libreoffice_timeout=300,
        ),
    )
    if render_result.page_count > max_pages:
        # Rendering finishes before the library reports the final page count.
        for image in render_result.images:
            image.path.unlink(missing_ok=True)
        raise DocumentPageLimitError(f"Document exceeds the {max_pages}-page limit")
    artifacts = build_artifacts(render_result.images, output_dir, text_blocks)
    manifest = {
        "schema_version": 1,
        "converter_version": "2026.09.0",
        "source": {"media_type": media_type_for(extension), "sha256": sha256_file(source)},
        "documents": [{"name": source.name, "pages": [asdict(item) for item in artifacts]}],
        "warnings": [],
    }
    (output_dir.parent / "manifest.json").write_text(
        json.dumps(manifest, ensure_ascii=False, indent=2), encoding="utf-8"
    )
    logger.info("Document rendering completed: filename=%s pages=%d", source.name, len(artifacts))
    return ConversionResult(artifacts, [])


def build_artifacts(
    rendered_images: tuple[RenderedImage, ...],
    output_dir: Path,
    text_blocks: list[str],
) -> list[Artifact]:
    artifacts: list[Artifact] = []
    for index, rendered_image in enumerate(rendered_images, start=1):
        image_path = rendered_image.path
        page_text = text_blocks[index - 1] if index <= len(text_blocks) else ""
        text_path = output_dir / f"page-{index:04d}.txt"
        text_path.write_text(page_text, encoding="utf-8")
        artifacts.append(
            Artifact(
                page_number=rendered_image.page_number,
                image_path=image_path.name,
                text_path=text_path.name,
                width=rendered_image.width,
                height=rendered_image.height,
                media_type="image/png",
                sha256=sha256_file(image_path),
            )
        )
    return artifacts


def extract_text(source: Path, extension: str) -> list[str]:
    if extension == ".pdf":
        import pymupdf

        with pymupdf.open(source) as document:
            return [page.get_text() for page in document]
    if extension == ".pptx":
        from pptx import Presentation

        presentation = Presentation(source)
        return [
            "\n".join(shape.text for shape in slide.shapes if getattr(shape, "has_text_frame", False))
            for slide in presentation.slides
        ]
    if extension == ".docx":
        from docx import Document

        document = Document(source)
        blocks = [paragraph.text for paragraph in document.paragraphs if paragraph.text]
        blocks.extend("\t".join(cell.text for cell in row.cells) for table in document.tables for row in table.rows)
        return ["\n".join(blocks)]
    from openpyxl import load_workbook

    workbook = load_workbook(source, data_only=False, read_only=True)
    try:
        return [
            f"<sheet name={json.dumps(sheet.title)}>\n"
            + "\n".join("\t".join("" if value is None else str(value) for value in row) for row in sheet.iter_rows(values_only=True))
            + "\n</sheet>"
            for sheet in workbook.worksheets
        ]
    finally:
        workbook.close()


def media_type_for(extension: str) -> str:
    return {
        ".pdf": "application/pdf",
        ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    }[extension]


def sha256_file(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as stream:
        for chunk in iter(lambda: stream.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()