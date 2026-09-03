"""Core document-to-image conversion services."""

from __future__ import annotations

import shutil
import subprocess
import sys
import tempfile
import zipfile
from dataclasses import dataclass
from pathlib import Path
from xml.etree import ElementTree

from PIL import Image


SUPPORTED_EXTENSIONS = {".pdf", ".pptx", ".docx", ".xlsx"}
OFFICE_PDF_FILTERS = {
    ".pptx": "impress_pdf_Export",
    ".docx": "writer_pdf_Export",
    ".xlsx": "calc_pdf_Export",
}
DRAWING_NAMESPACE = "http://schemas.openxmlformats.org/drawingml/2006/main"
PRESENTATION_NAMESPACE = "http://schemas.openxmlformats.org/presentationml/2006/main"


@dataclass(frozen=True)
class ConversionOptions:
    output_dir: Path
    image_format: str = "png"
    max_dimension: int = 2048
    pdf_dpi: int = 150


class ConversionError(RuntimeError):
    """Raised when an input document cannot be converted."""


def convert(input_path: Path, options: ConversionOptions) -> list[Path]:
    extension = input_path.suffix.lower()
    if extension not in SUPPORTED_EXTENSIONS:
        supported = ", ".join(sorted(SUPPORTED_EXTENSIONS))
        raise ConversionError(
            f"未対応の形式です: {extension or '(拡張子なし)'}。対応形式: {supported}"
        )

    options.output_dir.mkdir(parents=True, exist_ok=True)
    if extension == ".pdf":
        return convert_pdf(input_path, options)
    return convert_office(input_path, options, OFFICE_PDF_FILTERS[extension])


def convert_pdf(input_path: Path, options: ConversionOptions) -> list[Path]:
    try:
        import pymupdf
    except ImportError as error:
        raise ConversionError(
            "PyMuPDFが必要です。requirements.txtをインストールしてください。"
        ) from error

    try:
        document = pymupdf.open(input_path)
    except Exception as error:
        raise ConversionError(f"PDFを開けません: {error}") from error

    outputs: list[Path] = []
    with document:
        if document.needs_pass:
            raise ConversionError("暗号化されたPDFには対応していません。")

        matrix = pymupdf.Matrix(options.pdf_dpi / 72, options.pdf_dpi / 72)
        for page_number, page in enumerate(document, start=1):
            pixmap = page.get_pixmap(matrix=matrix, alpha=False)
            image = Image.frombytes("RGB", (pixmap.width, pixmap.height), pixmap.samples)
            outputs.append(save_image(image, f"page-{page_number}", page_number, options))
    return outputs


def convert_office(
    input_path: Path,
    options: ConversionOptions,
    pdf_filter: str,
) -> list[Path]:
    soffice = shutil.which("soffice")
    if soffice is None:
        raise ConversionError(
            "LibreOfficeが見つかりません。devcontainer内で実行してください。"
        )

    if input_path.suffix.lower() == ".pptx":
        warn_about_pptx_font_substitutions(input_path)

    with tempfile.TemporaryDirectory(prefix="office-convert-") as temporary_directory:
        work_dir = Path(temporary_directory)
        pdf_dir = work_dir / "pdf"
        profile_dir = work_dir / "profile"
        pdf_dir.mkdir()
        profile_dir.mkdir()

        office_input = input_path
        if input_path.suffix.lower() == ".pptx":
            office_input = work_dir / input_path.name
            prepare_pptx_for_libreoffice(input_path, office_input)

        command = build_soffice_command(
            soffice,
            office_input,
            pdf_dir,
            profile_dir,
            pdf_filter,
        )
        result = run_soffice(command, input_path.suffix.upper())
        converted_pdf = pdf_dir / f"{input_path.stem}.pdf"
        if result.returncode != 0 or not converted_pdf.is_file():
            details = (result.stderr or result.stdout).strip()
            message = (
                f"LibreOfficeによる{input_path.suffix.upper()}変換に失敗しました "
                f"(exit={result.returncode})"
            )
            if details:
                message += f": {details}"
            raise ConversionError(message)
        return convert_pdf(converted_pdf, options)


def prepare_pptx_for_libreoffice(input_path: Path, output_path: Path) -> None:
    with (
        zipfile.ZipFile(input_path) as source,
        zipfile.ZipFile(output_path, "w") as destination,
    ):
        presentation = ElementTree.fromstring(source.read("ppt/presentation.xml"))
        slide_size = presentation.find(f"{{{PRESENTATION_NAMESPACE}}}sldSz")
        dimensions = None
        if slide_size is not None:
            dimensions = (int(slide_size.get("cx", "0")), int(slide_size.get("cy", "0")))
        for item in source.infolist():
            content = source.read(item.filename)
            if item.filename.startswith("ppt/slides/slide") and item.filename.endswith(".xml"):
                content = normalize_pptx_slide(content, dimensions)
            destination.writestr(item, content)


def normalize_pptx_slide(
    content: bytes,
    slide_size: tuple[int, int] | None = None,
) -> bytes:
    root = ElementTree.fromstring(content)
    namespaces = {"a": DRAWING_NAMESPACE, "p": PRESENTATION_NAMESPACE}

    if slide_size is not None:
        normalize_radial_arrow_directions(root, namespaces, slide_size)

    for transform in root.findall(".//p:spPr/a:xfrm", namespaces):
        offset = transform.find("a:off", namespaces)
        extent = transform.find("a:ext", namespaces)
        if offset is None or extent is None:
            continue
        normalize_negative_extent(offset, extent, "x", "cx")
        normalize_negative_extent(offset, extent, "y", "cy")

    for shape in root.findall(".//p:sp", namespaces):
        prevent_unwanted_text_wrap(shape, namespaces)

    return ElementTree.tostring(root, encoding="utf-8", xml_declaration=True)


def normalize_negative_extent(
    offset: ElementTree.Element,
    extent: ElementTree.Element,
    offset_name: str,
    extent_name: str,
) -> None:
    dimension = int(extent.get(extent_name, "0"))
    if dimension >= 0:
        return
    offset.set(offset_name, str(int(offset.get(offset_name, "0")) + dimension))
    extent.set(extent_name, str(-dimension))


def prevent_unwanted_text_wrap(
    shape: ElementTree.Element,
    namespaces: dict[str, str],
) -> None:
    paragraphs = shape.findall("p:txBody/a:p", namespaces)
    if len(paragraphs) != 1:
        return
    text = "".join(node.text or "" for node in paragraphs[0].findall(".//a:t", namespaces))
    if not text:
        return
    transform = shape.find("p:spPr/a:xfrm/a:ext", namespaces)
    body_properties = shape.find("p:txBody/a:bodyPr", namespaces)
    run_properties = paragraphs[0].find(".//a:rPr", namespaces)
    if transform is None or body_properties is None or run_properties is None:
        return
    font_size = int(run_properties.get("sz", "0"))
    width = int(transform.get("cx", "0"))
    estimated_width = sum(1.0 if ord(character) > 0xFF else 0.55 for character in text)
    estimated_width *= font_size * 127
    if estimated_width > width:
        body_properties.set("wrap", "none")


def normalize_radial_arrow_directions(
    root: ElementTree.Element,
    namespaces: dict[str, str],
    slide_size: tuple[int, int],
) -> None:
    arrow_shapes = [
        shape
        for shape in root.findall(".//p:sp", namespaces)
        if shape.find("p:spPr/a:ln/a:tailEnd", namespaces) is not None
        and shape.find("p:spPr/a:prstGeom[@prst='line']", namespaces) is not None
    ]
    if len(arrow_shapes) < 3:
        return

    center_x, center_y = slide_size[0] / 2, slide_size[1] / 2
    for shape in arrow_shapes:
        transform = shape.find("p:spPr/a:xfrm", namespaces)
        offset = shape.find("p:spPr/a:xfrm/a:off", namespaces)
        extent = shape.find("p:spPr/a:xfrm/a:ext", namespaces)
        line = shape.find("p:spPr/a:ln", namespaces)
        if transform is None or offset is None or extent is None or line is None:
            continue

        x, y = int(offset.get("x", "0")), int(offset.get("y", "0"))
        width, height = int(extent.get("cx", "0")), int(extent.get("cy", "0"))
        if width < 0 or height < 0:
            continue
        start = (
            x + width if transform.get("flipH") in {"1", "true"} else x,
            y + height if transform.get("flipV") in {"1", "true"} else y,
        )
        end = (
            x if transform.get("flipH") in {"1", "true"} else x + width,
            y if transform.get("flipV") in {"1", "true"} else y + height,
        )
        start_distance = (start[0] - center_x) ** 2 + (start[1] - center_y) ** 2
        end_distance = (end[0] - center_x) ** 2 + (end[1] - center_y) ** 2
        if end_distance >= start_distance:
            continue

        tail = line.find("a:tailEnd", namespaces)
        if tail is not None:
            tail.tag = f"{{{DRAWING_NAMESPACE}}}headEnd"


def build_soffice_command(
    soffice: str,
    input_path: Path,
    pdf_dir: Path,
    profile_dir: Path,
    pdf_filter: str,
) -> list[str]:
    return [
        soffice,
        "--headless",
        "--nologo",
        "--nodefault",
        "--nofirststartwizard",
        f"-env:UserInstallation={profile_dir.as_uri()}",
        "--convert-to",
        f"pdf:{pdf_filter}",
        "--outdir",
        str(pdf_dir),
        str(input_path),
    ]


def run_soffice(command: list[str], extension: str) -> subprocess.CompletedProcess[str]:
    try:
        return subprocess.run(
            command,
            check=False,
            capture_output=True,
            text=True,
            timeout=60,
        )
    except subprocess.TimeoutExpired as error:
        raise ConversionError(
            f"LibreOfficeによる{extension}変換が60秒でタイムアウトしました。"
        ) from error


def save_image(
    image: Image.Image,
    stem: str,
    sequence: int,
    options: ConversionOptions,
) -> Path:
    extension = ".webp" if options.image_format == "webp" else ".png"
    output_path = options.output_dir / f"{sequence:04d}-{stem}{extension}"
    converted = fit_within(image.convert("RGB"), options.max_dimension)
    save_options: dict[str, object] = {"optimize": True}
    if options.image_format == "webp":
        save_options.update(lossless=True, method=6)
    converted.save(output_path, **save_options)
    return output_path


def fit_within(image: Image.Image, max_dimension: int) -> Image.Image:
    if max(image.size) <= max_dimension:
        return image
    resized = image.copy()
    resized.thumbnail((max_dimension, max_dimension), Image.Resampling.LANCZOS)
    return resized


def warn_about_pptx_font_substitutions(input_path: Path) -> None:
    fc_match = shutil.which("fc-match")
    if fc_match is None:
        return

    for requested_font in sorted(collect_pptx_fonts(input_path)):
        resolved_font = resolve_font_family(fc_match, requested_font)
        if resolved_font and resolved_font.casefold() != requested_font.casefold():
            print(
                f"warning: PPTX font substitution: {requested_font} -> {resolved_font}",
                file=sys.stderr,
            )


def collect_pptx_fonts(input_path: Path) -> set[str]:
    try:
        from pptx import Presentation

        presentation = Presentation(input_path)
        fonts = {
            run.font.name
            for slide in presentation.slides
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False)
            for paragraph in shape.text_frame.paragraphs
            for run in paragraph.runs
            if run.font.name
        }
        fonts.update(collect_pptx_theme_fonts(input_path))
        return fonts
    except Exception:
        return set()


def collect_pptx_theme_fonts(input_path: Path) -> set[str]:
    fonts: set[str] = set()
    with zipfile.ZipFile(input_path) as archive:
        theme_names = (
            name
            for name in archive.namelist()
            if name.startswith("ppt/theme/") and name.endswith(".xml")
        )
        for name in theme_names:
            root = ElementTree.fromstring(archive.read(name))
            for element in root.iter():
                local_name = element.tag.rsplit("}", 1)[-1]
                typeface = element.get("typeface")
                if local_name in {"latin", "ea", "cs"} and typeface and not typeface.startswith("+"):
                    fonts.add(typeface)
    return fonts


def resolve_font_family(fc_match: str, requested_font: str) -> str | None:
    try:
        result = subprocess.run(
            [fc_match, "--format=%{family[0]}", requested_font],
            check=False,
            capture_output=True,
            text=True,
            timeout=5,
        )
    except (OSError, subprocess.TimeoutExpired):
        return None
    if result.returncode != 0:
        return None
    return result.stdout.strip() or None