from __future__ import annotations

import tempfile
import unittest
from pathlib import Path
from unittest.mock import patch
from xml.etree import ElementTree

import pymupdf
from docx import Document
from openpyxl import Workbook
from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from convert_to_images import ConversionOptions, convert, main
from converter_core import DRAWING_NAMESPACE, normalize_pptx_slide


class ConvertToImagesTest(unittest.TestCase):
    def test_pptx_compatibility_normalizes_negative_extents_and_overflow(self) -> None:
        slide_xml = b"""<?xml version="1.0" encoding="UTF-8"?>
        <p:sld xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
               xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
          <p:cSld><p:spTree><p:sp><p:spPr><a:xfrm><a:off x="100" y="200"/>
            <a:ext cx="-40" cy="-50"/></a:xfrm></p:spPr>
            <p:txBody><a:bodyPr wrap="square"/><a:p><a:r><a:rPr sz="100"/>
              <a:t>Overflowing text</a:t></a:r></a:p></p:txBody>
          </p:sp></p:spTree></p:cSld>
        </p:sld>"""

        normalized = ElementTree.fromstring(normalize_pptx_slide(slide_xml))
        namespace = {"a": DRAWING_NAMESPACE}
        transform = normalized.find(".//a:xfrm", namespace)

        self.assertIsNotNone(transform)
        self.assertIsNone(transform.get("flipH"))
        self.assertIsNone(transform.get("flipV"))
        self.assertEqual(
            transform.find("a:off", namespace).attrib,
            {"x": "60", "y": "150"},
        )
        self.assertEqual(
            transform.find("a:ext", namespace).attrib,
            {"cx": "40", "cy": "50"},
        )
        self.assertEqual(
            normalized.find(".//a:bodyPr", namespace).get("wrap"),
            "none",
        )

    def test_pptx_compatibility_points_radial_arrows_away_from_center(self) -> None:
        arrows = "".join(
            f"""<p:sp><p:spPr><a:xfrm><a:off x="{x}" y="50"/>
            <a:ext cx="20" cy="0"/></a:xfrm><a:prstGeom prst="line"/>
            <a:ln><a:tailEnd type="triangle"/></a:ln></p:spPr></p:sp>"""
            for x in (10, 30, 70, 90)
        )
        slide_xml = f"""<p:sld
          xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
          xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
          <p:cSld><p:spTree>{arrows}</p:spTree></p:cSld></p:sld>""".encode()

        normalized = ElementTree.fromstring(normalize_pptx_slide(slide_xml, (100, 100)))
        namespace = {"a": DRAWING_NAMESPACE}
        lines = normalized.findall(".//a:ln", namespace)

        self.assertIsNotNone(lines[0].find("a:headEnd", namespace))
        self.assertIsNotNone(lines[1].find("a:headEnd", namespace))
        self.assertIsNotNone(lines[2].find("a:tailEnd", namespace))
        self.assertIsNotNone(lines[3].find("a:tailEnd", namespace))

    def test_libreoffice_renderer_uses_format_specific_pdf_filter(self) -> None:
        with tempfile.TemporaryDirectory() as temporary_directory:
            root = Path(temporary_directory)
            inputs = self._create_inputs(root)
            filters = {
                ".pptx": "impress_pdf_Export",
                ".docx": "writer_pdf_Export",
                ".xlsx": "calc_pdf_Export",
            }

            def create_converted_pdf(command: list[str], **_: object) -> object:
                output_dir = Path(command[command.index("--outdir") + 1])
                input_path = Path(command[-1])
                pdf = pymupdf.open()
                pdf.new_page()
                pdf.save(output_dir / f"{input_path.stem}.pdf")
                pdf.close()
                return type("Result", (), {"returncode": 0, "stdout": "", "stderr": ""})()

            for input_path in inputs[1:]:
                with self.subTest(extension=input_path.suffix):
                    options = self._options(root / f"output-{input_path.suffix[1:]}")
                    with (
                        patch("converter_core.shutil.which", return_value="/usr/bin/soffice"),
                        patch(
                            "converter_core.warn_about_pptx_font_substitutions"
                        ),
                        patch(
                            "converter_core.subprocess.run",
                            side_effect=create_converted_pdf,
                        ) as run,
                    ):
                        outputs = convert(input_path, options)

                    self.assertEqual(len(outputs), 1)
                    self.assertIn(f"pdf:{filters[input_path.suffix]}", run.call_args.args[0])

    def test_cli_accepts_input_file_argument(self) -> None:
        with tempfile.TemporaryDirectory() as temporary_directory:
            root = Path(temporary_directory)
            pdf_path = self._create_inputs(root)[0]
            output_dir = root / "cli-output"

            exit_code = main(
                [
                    str(pdf_path),
                    "--output-dir",
                    str(output_dir),
                    "--max-dimension",
                    "1024",
                ]
            )

            self.assertEqual(exit_code, 0)
            self.assertGreaterEqual(len(list(output_dir.glob("*.png"))), 1)

    def test_pdf_produces_readable_png_image(self) -> None:
        with tempfile.TemporaryDirectory() as temporary_directory:
            root = Path(temporary_directory)
            pdf_path = self._create_inputs(root)[0]
            outputs = convert(pdf_path, self._options(root / "output-pdf"))

            self.assertEqual(len(outputs), 1)
            with Image.open(outputs[0]) as image:
                image.verify()

    def _options(self, output_dir: Path) -> ConversionOptions:
        return ConversionOptions(
            output_dir=output_dir,
            image_format="png",
            max_dimension=1024,
            pdf_dpi=96,
        )

    def _create_inputs(self, root: Path) -> list[Path]:
        pdf_path = root / "sample.pdf"
        pdf = pymupdf.open()
        page = pdf.new_page()
        page.insert_text((72, 72), "PDF sample")
        pdf.save(pdf_path)
        pdf.close()

        pptx_path = root / "sample.pptx"
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        text_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
        text_box.text = "PPTX sample"
        presentation.save(pptx_path)

        docx_path = root / "sample.docx"
        document = Document()
        document.add_heading("DOCX sample", level=1)
        document.add_paragraph("This paragraph is rendered as a logical block.")
        embedded_image_path = root / "embedded.png"
        Image.new("RGB", (80, 40), "#2563eb").save(embedded_image_path)
        document.add_picture(str(embedded_image_path), width=Inches(1))
        table = document.add_table(rows=2, cols=2)
        table.cell(0, 0).text = "A"
        table.cell(0, 1).text = "B"
        table.cell(1, 0).text = "1"
        table.cell(1, 1).text = "2"
        document.save(docx_path)

        xlsx_path = root / "sample.xlsx"
        workbook = Workbook()
        sheet = workbook.active
        sheet.title = "Summary"
        sheet.append(["Item", "Value", "Formula"])
        sheet.append(["Alpha", 10, "=B2*2"])
        sheet.append(["Beta", 20, "=B3*2"])
        workbook.save(xlsx_path)
        workbook.close()

        return [pdf_path, pptx_path, docx_path, xlsx_path]


if __name__ == "__main__":
    unittest.main()