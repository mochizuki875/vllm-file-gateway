from __future__ import annotations

import hashlib
import json
import shutil
import subprocess
import tempfile
import zipfile
from dataclasses import asdict, dataclass
from pathlib import Path
from xml.etree import ElementTree

from PIL import Image


SUPPORTED_EXTENSIONS = {".pdf", ".pptx", ".docx", ".xlsx"}
OFFICE_PDF_FILTERS = {
    ".pptx": "impress_pdf_Export",
    ".docx": "writer_pdf_Export",
    ".xlsx": "calc_pdf_Export",
}
DRAWING_NAMESPACE = "http://schemas.openxmlformats.org/drawingml/2006/main"
PRESENTATION_NAMESPACE = "http://schemas.openxmlformats.org/presentationml/2006/main"


class ConversionError(RuntimeError):
    pass


@dataclass(frozen=True)
class Artifact:
    page_number: int
    image_path: str
    text_path: str
    width: int
    height: int
    media_type: str
    sha256: str


@dataclass(frozen=True)
class ConversionResult:
    artifacts: list[Artifact]
    warnings: list[str]


def convert_document(source: Path, output_dir: Path, max_dimension: int = 2048) -> ConversionResult:
    extension = source.suffix.lower()
    if extension not in SUPPORTED_EXTENSIONS:
        raise ConversionError(f"Unsupported file type: {extension or '(none)'}")
    output_dir.mkdir(parents=True, exist_ok=True)
    text_blocks = extract_text(source, extension)
    pdf_path = source
    temporary: tempfile.TemporaryDirectory[str] | None = None
    try:
        if extension != ".pdf":
            temporary = tempfile.TemporaryDirectory(prefix="gateway-office-")
            pdf_path = convert_office_to_pdf(source, Path(temporary.name), extension)
        artifacts = render_pdf(pdf_path, output_dir, text_blocks, max_dimension)
    finally:
        if temporary:
            temporary.cleanup()
    manifest = {
        "schema_version": 1,
        "converter_version": "2026.09.0",
        "source": {"media_type": media_type_for(extension), "sha256": sha256_file(source)},
        "documents": [{"name": source.name, "pages": [asdict(item) for item in artifacts]}],
        "warnings": [],
    }
    (output_dir.parent / "manifest.json").write_text(
        json.dumps(manifest, ensure_ascii=False, indent=2), encoding="utf-8"
    )
    return ConversionResult(artifacts, [])


def render_pdf(pdf_path: Path, output_dir: Path, text_blocks: list[str], max_dimension: int) -> list[Artifact]:
    import pymupdf

    try:
        document = pymupdf.open(pdf_path)
    except Exception as error:
        raise ConversionError(f"Unable to open PDF: {error}") from error
    artifacts: list[Artifact] = []
    with document:
        if document.needs_pass:
            raise ConversionError("Encrypted PDFs are not supported")
        if len(document) > 20:
            raise ConversionError("Document exceeds the 20 page/image limit")
        matrix = pymupdf.Matrix(150 / 72, 150 / 72)
        for index, page in enumerate(document, start=1):
            pixmap = page.get_pixmap(matrix=matrix, alpha=False)
            image = Image.frombytes("RGB", (pixmap.width, pixmap.height), pixmap.samples)
            image.thumbnail((max_dimension, max_dimension), Image.Resampling.LANCZOS)
            image_path = output_dir / f"page-{index:04d}.webp"
            image.save(image_path, "WEBP", lossless=True, method=6)
            page_text = text_blocks[index - 1] if index <= len(text_blocks) else page.get_text()
            text_path = output_dir / f"page-{index:04d}.txt"
            text_path.write_text(page_text, encoding="utf-8")
            artifacts.append(
                Artifact(
                    page_number=index,
                    image_path=image_path.name,
                    text_path=text_path.name,
                    width=image.width,
                    height=image.height,
                    media_type="image/webp",
                    sha256=sha256_file(image_path),
                )
            )
    return artifacts


def extract_text(source: Path, extension: str) -> list[str]:
    if extension == ".pdf":
        import pymupdf

        with pymupdf.open(source) as document:
            return [page.get_text() for page in document]
    if extension == ".pptx":
        from pptx import Presentation

        presentation = Presentation(source)
        return [
            "\n".join(shape.text for shape in slide.shapes if getattr(shape, "has_text_frame", False))
            for slide in presentation.slides
        ]
    if extension == ".docx":
        from docx import Document

        document = Document(source)
        blocks = [paragraph.text for paragraph in document.paragraphs if paragraph.text]
        blocks.extend("\t".join(cell.text for cell in row.cells) for table in document.tables for row in table.rows)
        return ["\n".join(blocks)]
    from openpyxl import load_workbook

    workbook = load_workbook(source, data_only=False, read_only=True)
    try:
        return [
            f"<sheet name={json.dumps(sheet.title)}>\n"
            + "\n".join("\t".join("" if value is None else str(value) for value in row) for row in sheet.iter_rows(values_only=True))
            + "\n</sheet>"
            for sheet in workbook.worksheets
        ]
    finally:
        workbook.close()


def convert_office_to_pdf(source: Path, work_dir: Path, extension: str) -> Path:
    soffice = shutil.which("soffice")
    if not soffice:
        raise ConversionError("LibreOffice (soffice) is required")
    pdf_dir = work_dir / "pdf"
    profile_dir = work_dir / "profile"
    pdf_dir.mkdir(parents=True)
    profile_dir.mkdir()
    office_input = source
    if extension == ".pptx":
        office_input = work_dir / source.name
        prepare_pptx_for_libreoffice(source, office_input)
    command = [
        soffice,
        "--headless",
        "--nologo",
        "--nodefault",
        "--nofirststartwizard",
        f"-env:UserInstallation={profile_dir.as_uri()}",
        "--convert-to",
        f"pdf:{OFFICE_PDF_FILTERS[extension]}",
        "--outdir",
        str(pdf_dir),
        str(office_input),
    ]
    try:
        result = subprocess.run(command, check=False, capture_output=True, text=True, timeout=300)
    except subprocess.TimeoutExpired as error:
        raise ConversionError("LibreOffice conversion timed out") from error
    converted = pdf_dir / f"{source.stem}.pdf"
    if result.returncode != 0 or not converted.is_file():
        details = (result.stderr or result.stdout).strip()
        raise ConversionError(f"LibreOffice conversion failed: {details}")
    return converted


def prepare_pptx_for_libreoffice(source: Path, destination_path: Path) -> None:
    with zipfile.ZipFile(source) as archive, zipfile.ZipFile(destination_path, "w") as destination:
        presentation = ElementTree.fromstring(archive.read("ppt/presentation.xml"))
        slide_size = presentation.find(f"{{{PRESENTATION_NAMESPACE}}}sldSz")
        dimensions = None if slide_size is None else (int(slide_size.get("cx", "0")), int(slide_size.get("cy", "0")))
        for item in archive.infolist():
            content = archive.read(item.filename)
            if item.filename.startswith("ppt/slides/slide") and item.filename.endswith(".xml"):
                content = normalize_pptx_slide(content, dimensions)
            destination.writestr(item, content)


def normalize_pptx_slide(content: bytes, slide_size: tuple[int, int] | None = None) -> bytes:
    root = ElementTree.fromstring(content)
    namespaces = {"a": DRAWING_NAMESPACE, "p": PRESENTATION_NAMESPACE}
    if slide_size:
        normalize_radial_arrows(root, namespaces, slide_size)
    for transform in root.findall(".//p:spPr/a:xfrm", namespaces):
        offset = transform.find("a:off", namespaces)
        extent = transform.find("a:ext", namespaces)
        if offset is None or extent is None:
            continue
        for offset_name, extent_name in (("x", "cx"), ("y", "cy")):
            value = int(extent.get(extent_name, "0"))
            if value < 0:
                offset.set(offset_name, str(int(offset.get(offset_name, "0")) + value))
                extent.set(extent_name, str(-value))
    for shape in root.findall(".//p:sp", namespaces):
        prevent_text_wrap(shape, namespaces)
    return ElementTree.tostring(root, encoding="utf-8", xml_declaration=True)


def prevent_text_wrap(shape: ElementTree.Element, namespaces: dict[str, str]) -> None:
    paragraphs = shape.findall("p:txBody/a:p", namespaces)
    if len(paragraphs) != 1:
        return
    text = "".join(node.text or "" for node in paragraphs[0].findall(".//a:t", namespaces))
    extent = shape.find("p:spPr/a:xfrm/a:ext", namespaces)
    body = shape.find("p:txBody/a:bodyPr", namespaces)
    run = paragraphs[0].find(".//a:rPr", namespaces)
    if text and extent is not None and body is not None and run is not None:
        estimated = sum(1.0 if ord(character) > 0xFF else 0.55 for character in text) * int(run.get("sz", "0")) * 127
        if estimated > int(extent.get("cx", "0")):
            body.set("wrap", "none")


def normalize_radial_arrows(root: ElementTree.Element, namespaces: dict[str, str], slide_size: tuple[int, int]) -> None:
    shapes = [shape for shape in root.findall(".//p:sp", namespaces) if shape.find("p:spPr/a:ln/a:tailEnd", namespaces) is not None and shape.find("p:spPr/a:prstGeom[@prst='line']", namespaces) is not None]
    if len(shapes) < 3:
        return
    center = (slide_size[0] / 2, slide_size[1] / 2)
    for shape in shapes:
        transform = shape.find("p:spPr/a:xfrm", namespaces)
        offset = shape.find("p:spPr/a:xfrm/a:off", namespaces)
        extent = shape.find("p:spPr/a:xfrm/a:ext", namespaces)
        line = shape.find("p:spPr/a:ln", namespaces)
        if transform is None or offset is None or extent is None or line is None:
            continue
        x, y = int(offset.get("x", "0")), int(offset.get("y", "0"))
        width, height = int(extent.get("cx", "0")), int(extent.get("cy", "0"))
        if width < 0 or height < 0:
            continue
        start = (x + width if transform.get("flipH") in {"1", "true"} else x, y + height if transform.get("flipV") in {"1", "true"} else y)
        end = (x if transform.get("flipH") in {"1", "true"} else x + width, y if transform.get("flipV") in {"1", "true"} else y + height)
        distance = lambda point: (point[0] - center[0]) ** 2 + (point[1] - center[1]) ** 2
        if distance(end) < distance(start):
            tail = line.find("a:tailEnd", namespaces)
            if tail is not None:
                tail.tag = f"{{{DRAWING_NAMESPACE}}}headEnd"


def media_type_for(extension: str) -> str:
    return {
        ".pdf": "application/pdf",
        ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    }[extension]


def sha256_file(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as stream:
        for chunk in iter(lambda: stream.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()